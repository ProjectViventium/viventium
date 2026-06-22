#!/usr/bin/env python3
"""Run a public-safe local GlassHive UI fixture for browser QA.

This script starts the real FastAPI app with an injected synthetic runtime,
creates one project, one Codex-profile worker, one completed run, and a small
artifact set. It is intentionally provider-free and cloud-free so user-grade UI
checks can run before any live integration QA.
"""

from __future__ import annotations

import argparse
import json
import os
import shutil
import signal
import socket
import sys
import tempfile
import textwrap
import threading
import time
import zipfile
from pathlib import Path
from typing import Any

import uvicorn
from fastapi.responses import HTMLResponse
from fastapi.testclient import TestClient


REPO_ROOT = Path(__file__).resolve().parents[3]
RUNTIME_SRC = REPO_ROOT / "viventium_v0_4" / "GlassHive" / "runtime_phase1" / "src"
sys.path.insert(0, str(RUNTIME_SRC))

from workers_projects_runtime.api import create_app  # noqa: E402
from workers_projects_runtime.openclaw_runtime import RuntimeInfo  # noqa: E402
from workers_projects_runtime.terminal_takeover import TerminalTarget  # noqa: E402


PUBLIC_MARKER = "GLASSHIVE_BROWSER_QA_MARKER"


def _free_port() -> int:
    with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as sock:
        sock.bind(("127.0.0.1", 0))
        return int(sock.getsockname()[1])


def _write_zip(path: Path, members: dict[str, str | bytes]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for name, value in members.items():
            data = value.encode("utf-8") if isinstance(value, str) else value
            archive.writestr(name, data)


def _write_minimal_pdf(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    content = b"BT /F1 18 Tf 72 720 Td (GlassHive browser QA fixture) Tj ET"
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length %d >>\nstream\n%s\nendstream" % (len(content), content),
    ]
    output = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for index, obj in enumerate(objects, start=1):
        offsets.append(len(output))
        output.extend(f"{index} 0 obj\n".encode("ascii"))
        output.extend(obj)
        output.extend(b"\nendobj\n")
    xref = len(output)
    output.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    output.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        output.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    output.extend(f"trailer << /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n".encode("ascii"))
    path.write_bytes(bytes(output))


def _write_docx(path: Path) -> None:
    try:
        from docx import Document

        document = Document()
        document.add_heading("GlassHive Browser QA Fixture", level=1)
        document.add_paragraph(f"This document contains {PUBLIC_MARKER}.")
        path.parent.mkdir(parents=True, exist_ok=True)
        document.save(path)
        return
    except Exception:
        pass
    _write_zip(
        path,
        {
            "[Content_Types].xml": """<?xml version="1.0" encoding="UTF-8"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>""",
            "_rels/.rels": """<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>""",
            "word/document.xml": f"""<?xml version="1.0" encoding="UTF-8"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:body><w:p><w:r><w:t>{PUBLIC_MARKER}</w:t></w:r></w:p></w:body>
</w:document>""",
        },
    )


def _write_xlsx(path: Path) -> None:
    try:
        from openpyxl import Workbook

        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "QA"
        sheet.append(["marker", "status"])
        sheet.append([PUBLIC_MARKER, "pass"])
        path.parent.mkdir(parents=True, exist_ok=True)
        workbook.save(path)
        return
    except Exception:
        pass
    _write_zip(
        path,
        {
            "[Content_Types].xml": """<?xml version="1.0" encoding="UTF-8"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/xl/workbook.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/>
</Types>""",
            "_rels/.rels": """<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="xl/workbook.xml"/>
</Relationships>""",
            "xl/workbook.xml": """<?xml version="1.0" encoding="UTF-8"?>
<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main"><sheets/></workbook>""",
        },
    )


def _write_pptx(path: Path) -> None:
    try:
        from pptx import Presentation

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "GlassHive Browser QA Fixture"
        textbox = slide.shapes.add_textbox(914400, 1828800, 7315200, 914400)
        textbox.text_frame.text = PUBLIC_MARKER
        path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(path)
        return
    except Exception:
        pass
    _write_zip(
        path,
        {
            "[Content_Types].xml": """<?xml version="1.0" encoding="UTF-8"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>
</Types>""",
            "_rels/.rels": """<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>
</Relationships>""",
            "ppt/presentation.xml": f"""<?xml version="1.0" encoding="UTF-8"?>
<p:presentation xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:sldIdLst/><p:notesSz cx="6858000" cy="9144000"/><p:embeddedFontLst/>
</p:presentation><!-- {PUBLIC_MARKER} -->""",
        },
    )


def _write_artifacts(workspace: Path) -> None:
    workspace.mkdir(parents=True, exist_ok=True)
    (workspace / "answer.md").write_text(
        textwrap.dedent(
            f"""\
            # GlassHive Browser QA Fixture

            {PUBLIC_MARKER}

            This synthetic deliverable proves the local UI can render text artifacts,
            expose short-link actions, and preserve workspace navigation without raw
            signed-token leakage.
            """
        ),
        encoding="utf-8",
    )
    output_dir = workspace / "output"
    reports_dir = workspace / "reports"
    artifacts_dir = workspace / "artifacts"
    output_dir.mkdir(parents=True, exist_ok=True)
    reports_dir.mkdir(parents=True, exist_ok=True)
    artifacts_dir.mkdir(parents=True, exist_ok=True)
    (output_dir / "data.csv").write_text("marker,status\nGLASSHIVE_BROWSER_QA_MARKER,pass\n", encoding="utf-8")
    (reports_dir / "report.html").write_text(
        f"<!doctype html><title>GlassHive QA</title><main><h1>{PUBLIC_MARKER}</h1></main>",
        encoding="utf-8",
    )
    _write_minimal_pdf(artifacts_dir / "report.pdf")
    _write_xlsx(artifacts_dir / "book.xlsx")
    _write_docx(artifacts_dir / "brief.docx")
    _write_pptx(artifacts_dir / "deck.pptx")

    scratch = workspace / "tmp" / "chrome-user-data" / "Default"
    scratch.mkdir(parents=True, exist_ok=True)
    (scratch / "internal.html").write_text("<script>window.secret = true;</script>", encoding="utf-8")
    uploads = workspace / "uploads"
    uploads.mkdir(parents=True, exist_ok=True)
    (uploads / "source.txt.metadata.json").write_text('{"internal": true}', encoding="utf-8")


class SyntheticQaRuntime:
    def __init__(self, state_root: Path) -> None:
        self.state_root = state_root
        self.runtime_root = state_root / "runtime"
        self.workers_root = self.runtime_root / "workers"
        self.state_root.mkdir(parents=True, exist_ok=True)

    def _paths(self, worker: dict[str, Any]) -> tuple[Path, Path]:
        worker_id = str(worker["worker_id"])
        state_dir = self.workers_root / worker_id / "state"
        workspace_dir = self.workers_root / worker_id / "workspace"
        state_dir.mkdir(parents=True, exist_ok=True)
        workspace_dir.mkdir(parents=True, exist_ok=True)
        return state_dir, workspace_dir

    def resolve_model(self, profile: str) -> str:
        return f"qa-synthetic/{profile or 'default'}"

    def preflight_worker_profile(self, profile: str, execution_mode: str = "docker") -> None:
        _ = profile, execution_mode

    def ensure_worker_ready(self, worker: dict[str, Any]) -> RuntimeInfo:
        state_dir, workspace_dir = self._paths(worker)
        return RuntimeInfo(
            runtime=str(worker.get("profile") or "codex-cli"),
            model=str(worker.get("model") or self.resolve_model(str(worker.get("profile") or "codex-cli"))),
            gateway_url=f"http://127.0.0.1/synthetic/{worker['worker_id']}",
            gateway_port=None,
            gateway_token=None,
            session_key=f"qa:{worker['worker_id']}",
            state_dir=str(state_dir),
            workspace_dir=str(workspace_dir),
            pid=os.getpid(),
        )

    def pause_worker(self, worker: dict[str, Any]) -> RuntimeInfo:
        return self.ensure_worker_ready(worker)

    def interrupt_worker(self, worker: dict[str, Any], run_id: str | None = None) -> RuntimeInfo:
        _ = run_id
        return self.ensure_worker_ready(worker)

    def terminate_worker(self, worker: dict[str, Any]) -> RuntimeInfo:
        return self.ensure_worker_ready(worker)

    def reconcile_worker(self, worker: dict[str, Any]) -> RuntimeInfo:
        return self.ensure_worker_ready(worker)

    def run_task(
        self,
        worker: dict[str, Any],
        instruction: str,
        timeout_sec: float | None = None,
        run_id: str | None = None,
    ) -> str:
        _ = timeout_sec, run_id
        info = self.ensure_worker_ready(worker)
        workspace = Path(str(info.workspace_dir))
        if "GH_QA_SLEEP" in instruction:
            time.sleep(30)
        _write_artifacts(workspace)
        transcript = Path(str(info.state_dir)) / "synthetic-transcript.txt"
        transcript.write_text(f"Instruction received:\n{instruction}\n\n{PUBLIC_MARKER}\n", encoding="utf-8")
        return (
            "FINAL REPORT:\n"
            f"- Created public-safe fixture artifacts containing {PUBLIC_MARKER}.\n"
            "- Verified synthetic file input/output path for local browser QA.\n"
            "- No provider, cloud, LibreChat, or client environment was touched.\n"
        )

    def describe_worker(self, worker: dict[str, Any]) -> dict[str, Any]:
        info = self.ensure_worker_ready(worker)
        return {
            "mode": "workstation-desktop",
            "runtime": str(info.runtime),
            "sandbox_state": "ready",
            "workspace_dir": str(info.workspace_dir),
            "state_dir": str(info.state_dir),
            "view_url": f"/synthetic-desktop/{worker['worker_id']}",
        }

    def desktop_action(
        self,
        worker: dict[str, Any],
        action: str,
        url: str | None = None,
        run_id: str | None = None,
    ) -> dict[str, Any]:
        _ = url, run_id
        return {
            "action": action,
            "status": "launched",
            "mode": "workstation-desktop",
            "view_url": f"/synthetic-desktop/{worker['worker_id']}",
            "notes": "Synthetic local QA runtime.",
        }

    def terminal_target(self, worker: dict[str, Any]) -> TerminalTarget:
        info = self.ensure_worker_ready(worker)
        return TerminalTarget(
            command=[sys.executable, "-i"],
            cwd=str(info.workspace_dir),
            env={"TERM": "xterm-256color"},
            title="Synthetic GlassHive terminal",
            subtitle="Local browser QA fixture",
        )


def _poll_run(client: TestClient, run_id: str) -> dict[str, Any]:
    deadline = time.time() + 15
    last: dict[str, Any] = {}
    while time.time() < deadline:
        response = client.get(f"/v1/runs/{run_id}")
        response.raise_for_status()
        last = response.json()
        if last.get("state") in {"completed", "failed", "cancelled", "interrupted", "paused"}:
            return last
        time.sleep(0.1)
    raise TimeoutError(f"Run did not finish in fixture: {last}")


def build_fixture(state_root: Path) -> tuple[Any, dict[str, Any]]:
    os.environ.setdefault("VIVENTIUM_DISABLE_DEFAULT_RUNTIME_ENV", "1")
    os.environ["GLASSHIVE_SIGNED_LINK_SECRET"] = "public-safe-signed-link-secret"
    os.environ["GLASSHIVE_LINK_REF_STATE_PATH"] = str(state_root / "link_refs.sqlite3")
    os.environ.setdefault("GLASSHIVE_LINK_REF_TTL_SECONDS", "0")
    os.environ.setdefault("GLASSHIVE_DEFAULT_WORKER_PROFILE", "codex-cli")
    os.environ.setdefault("GLASSHIVE_ALLOWED_WORKER_PROFILES", "codex-cli,claude-code,openclaw-general")
    os.environ.setdefault("WPR_DEFAULT_EXECUTION_MODE", "docker")

    db_path = state_root / "runtime_phase1.sqlite3"
    app = create_app(db_path=str(db_path), runtime_backend="stub", runtime=SyntheticQaRuntime(state_root))

    @app.get("/synthetic-desktop/{worker_id}", response_class=HTMLResponse)
    def synthetic_desktop(worker_id: str) -> str:
        return f"""<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>Synthetic GlassHive Desktop</title>
  <style>
    body {{ margin: 0; background: #111827; color: #f9fafb; font-family: system-ui, sans-serif; }}
    main {{ min-height: 100vh; display: grid; place-content: center; gap: 12px; text-align: center; }}
    code {{ background: #1f2937; border: 1px solid #374151; border-radius: 6px; padding: 4px 8px; }}
  </style>
</head>
<body>
  <main>
    <h1>Synthetic GlassHive Desktop</h1>
    <p>{PUBLIC_MARKER}</p>
    <code>managed workspace</code>
  </main>
</body>
</html>"""

    client = TestClient(app)
    project = client.post(
        "/v1/projects",
        json={
            "owner_id": "public-safe-owner",
            "title": "Public-safe local browser QA",
            "goal": "Verify GlassHive UI artifacts, short links, and profile truth locally.",
            "default_worker_profile": "codex-cli",
        },
    )
    project.raise_for_status()
    project_json = project.json()
    worker = client.post(
        f"/v1/projects/{project_json['project_id']}/workers",
        json={
            "owner_id": "public-safe-owner",
            "name": "Codex fixture worker",
            "role": "Local UI QA worker",
            "profile": "codex-cli",
            "backend": "openclaw",
            "execution_mode": "docker",
            "start_synchronously": True,
        },
    )
    worker.raise_for_status()
    worker_json = worker.json()
    run = client.post(
        f"/v1/workers/{worker_json['worker_id']}/assign",
        json={
            "instruction": (
                "Create public-safe deliverables for local GlassHive browser QA. "
                "Include Markdown, CSV, HTML, PDF, XLSX, DOCX, and PPTX outputs."
            ),
            "effort": "xhigh",
        },
    )
    run.raise_for_status()
    run_json = _poll_run(client, run.json()["run_id"])
    artifacts = client.get(f"/v1/workers/{worker_json['worker_id']}/artifacts")
    artifacts.raise_for_status()
    live = client.get(f"/v1/workers/{worker_json['worker_id']}/live")
    live.raise_for_status()
    return app, {
        "project": project_json,
        "worker": worker_json,
        "run": run_json,
        "artifacts": artifacts.json().get("items", []),
        "live": live.json(),
        "db_path": str(db_path),
    }


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--state-root",
        type=Path,
        default=Path(tempfile.gettempdir()) / "glasshive-local-user-grade-fixture",
    )
    parser.add_argument("--port", type=int, default=0)
    parser.add_argument("--fresh", action="store_true", help="Delete the state root before starting.")
    args = parser.parse_args()

    if args.fresh and args.state_root.exists():
        shutil.rmtree(args.state_root)
    args.state_root.mkdir(parents=True, exist_ok=True)

    app, fixture = build_fixture(args.state_root)
    port = args.port or _free_port()
    base_url = f"http://127.0.0.1:{port}"
    project_id = fixture["project"]["project_id"]
    worker_id = fixture["worker"]["worker_id"]
    artifact_path = "answer.md"
    fixture.update(
        {
            "base_url": base_url,
            "project_url": f"{base_url}/ui/projects/{project_id}?worker_id={worker_id}",
            "worker_url": f"{base_url}/ui/workers/{worker_id}",
            "view_url": f"{base_url}/ui/workers/{worker_id}/view",
            "artifact_open_url": f"{base_url}/v1/workers/{worker_id}/artifacts/open?path={artifact_path}",
            "artifact_download_url": f"{base_url}/v1/workers/{worker_id}/artifacts/download?path={artifact_path}",
        }
    )

    stop_event = threading.Event()

    def _handle_stop(_signum: int, _frame: Any) -> None:
        stop_event.set()

    signal.signal(signal.SIGTERM, _handle_stop)
    signal.signal(signal.SIGINT, _handle_stop)

    config = uvicorn.Config(app, host="127.0.0.1", port=port, log_level="warning")
    server = uvicorn.Server(config)
    thread = threading.Thread(target=server.run, daemon=True)
    thread.start()
    deadline = time.time() + 10
    while not server.started and time.time() < deadline:
        time.sleep(0.05)
    if not server.started:
        raise RuntimeError("Fixture server did not start")

    print("GLASSHIVE_QA_SERVER " + json.dumps(fixture, sort_keys=True), flush=True)
    while not stop_event.is_set():
        time.sleep(0.2)
    server.should_exit = True
    thread.join(timeout=5)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
