#!/usr/bin/env python3
"""Run public-safe GlassHive deep-research/document-generation benchmark fixtures.

The suite is intentionally provider-free. It exercises the same constraint ledger,
artifact inventory, document validation, completion compliance, and content-hygiene
evidence paths used by real worker runs without encoding any private benchmark,
client, industry, or prompt-specific logic.
"""

from __future__ import annotations

import argparse
import csv
import json
import shutil
import sys
import tempfile
import textwrap
import time
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Callable


REPO_ROOT = Path(__file__).resolve().parents[3]
RUNTIME_SRC = REPO_ROOT / "viventium_v0_4" / "GlassHive" / "runtime_phase1" / "src"
sys.path.insert(0, str(RUNTIME_SRC))

from workers_projects_runtime.run_evidence import (  # noqa: E402
    build_constraint_ledger,
    build_run_evidence,
    write_constraint_ledger,
    write_run_evidence,
)


PUBLIC_MARKER = "GLASSHIVE_PUBLIC_BENCHMARK_MARKER"


@dataclass(frozen=True)
class BenchmarkCase:
    case_id: str
    description: str
    instruction: str
    expected_status: str
    expected_constraint_status: str
    expected_completion_status: str
    expected_coverage_status: str
    expected_hygiene_status: str
    setup: Callable[[Path], str]


def _write_csv(path: Path, rows: list[dict[str, str]]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", newline="", encoding="utf-8") as handle:
        writer = csv.DictWriter(handle, fieldnames=list(rows[0].keys()))
        writer.writeheader()
        writer.writerows(rows)


def _write_minimal_pdf(path: Path, title: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    content = f"BT /F1 18 Tf 72 720 Td ({title}) Tj ET".encode("utf-8")
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length %d >>\nstream\n" % len(content) + content + b"\nendstream",
    ]
    body = b"%PDF-1.4\n"
    offsets = [0]
    for index, obj in enumerate(objects, start=1):
        offsets.append(len(body))
        body += f"{index} 0 obj\n".encode("ascii") + obj + b"\nendobj\n"
    xref_offset = len(body)
    body += f"xref\n0 {len(objects) + 1}\n0000000000 65535 f \n".encode("ascii")
    for offset in offsets[1:]:
        body += f"{offset:010d} 00000 n \n".encode("ascii")
    body += (
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
        f"startxref\n{xref_offset}\n%%EOF\n"
    ).encode("ascii")
    path.write_bytes(body)


def _write_xlsx(path: Path, rows: list[list[str]]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    try:
        import openpyxl  # type: ignore

        workbook = openpyxl.Workbook()
        sheet = workbook.active
        sheet.title = "Evidence"
        for row in rows:
            sheet.append(row)
        workbook.save(path)
        return
    except Exception:
        pass
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", "<Types></Types>")
        archive.writestr("xl/workbook.xml", "<workbook></workbook>")


def _write_docx(path: Path, paragraphs: list[str]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    try:
        import docx  # type: ignore

        document = docx.Document()
        for paragraph in paragraphs:
            document.add_paragraph(paragraph)
        document.save(path)
        return
    except Exception:
        pass
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", "<Types></Types>")
        archive.writestr("word/document.xml", "<w:document></w:document>")


def _write_pptx(path: Path, title: str, body: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    try:
        import pptx  # type: ignore

        presentation = pptx.Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = body
        presentation.save(path)
        return
    except Exception:
        pass
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", "<Types></Types>")
        archive.writestr("ppt/presentation.xml", "<p:presentation></p:presentation>")


def _market_research_package(workspace: Path) -> str:
    report = workspace / "reports" / "market-research.md"
    report.parent.mkdir(parents=True, exist_ok=True)
    report.write_text(
        textwrap.dedent(
            f"""
            # Public-Safe Market Research Fixture

            {PUBLIC_MARKER}

            Scope preserved: only public/official source placeholders dated through June 2026.

            ## Seed Entities

            - Alpha Robotics: official product note, May 2026.
            - Beta Energy: official sustainability update, June 2026.

            ## Result

            Both seed entities are covered and no out-of-window source is used.
            """
        ).strip()
        + "\n",
        encoding="utf-8",
    )
    _write_csv(
        workspace / "output" / "market-screen.csv",
        [
            {"entity": "Alpha Robotics", "source_date": "May 2026", "fit": "covered"},
            {"entity": "Beta Energy", "source_date": "June 2026", "fit": "covered"},
        ],
    )
    _write_xlsx(
        workspace / "output" / "market-screen.xlsx",
        [
            ["entity", "source_date", "fit"],
            ["Alpha Robotics", "May 2026", "covered"],
            ["Beta Energy", "June 2026", "covered"],
        ],
    )
    return "FINAL REPORT:\nCreated Markdown, CSV, and XLSX deliverables for both seed entities."


def _technical_literature_brief(workspace: Path) -> str:
    report = workspace / "reports" / "technical-literature.md"
    report.parent.mkdir(parents=True, exist_ok=True)
    report.write_text(
        textwrap.dedent(
            f"""
            # Public-Safe Technical Literature Fixture

            {PUBLIC_MARKER}

            Source discipline: public technical literature and standards material dated through June 2026.

            - Container sandbox controls: public standard update, April 2026.
            - Local MCP transport auth: public protocol note, June 2026.

            The answer keeps both topics separate and records open risks without widening scope.
            """
        ).strip()
        + "\n",
        encoding="utf-8",
    )
    evidence = {
        "marker": PUBLIC_MARKER,
        "topics": [
            {"topic": "container sandbox controls", "source_date": "April 2026", "status": "covered"},
            {"topic": "local MCP transport auth", "source_date": "June 2026", "status": "covered"},
        ],
    }
    (workspace / "output").mkdir(parents=True, exist_ok=True)
    (workspace / "output" / "technical-evidence.json").write_text(json.dumps(evidence, indent=2) + "\n")
    return "FINAL REPORT:\nCreated Markdown and JSON evidence deliverables."


def _file_transform_package(workspace: Path) -> str:
    uploads = workspace / "uploads"
    uploads.mkdir(parents=True, exist_ok=True)
    _write_csv(
        uploads / "input.csv",
        [
            {"metric": "alpha", "value": "1"},
            {"metric": "beta", "value": "2"},
        ],
    )
    output = workspace / "output"
    _write_minimal_pdf(output / "transformed-report.pdf", "GlassHive file transform fixture")
    _write_docx(output / "transformed-brief.docx", [PUBLIC_MARKER, "The uploaded CSV was transformed into a brief."])
    _write_pptx(output / "transformed-deck.pptx", "Transform Fixture", PUBLIC_MARKER)
    (output / "transformed-summary.html").write_text(
        f"<!doctype html><html><head><title>Transform Fixture</title></head><body>{PUBLIC_MARKER}</body></html>\n",
        encoding="utf-8",
    )
    return "FINAL REPORT:\nCreated PDF, DOCX, PPTX, and HTML deliverables from the uploaded CSV."


def _text_only_answer(workspace: Path) -> str:
    (workspace / "work-log.md").write_text("No user-facing files requested.\n", encoding="utf-8")
    return f"FINAL REPORT:\n{PUBLIC_MARKER} Answered directly with no files or artifacts, as requested."


def _contaminated_structured_field(workspace: Path) -> str:
    _write_csv(
        workspace / "output" / "research-summary.csv",
        [
            {"entity": "Alpha Robotics", "notes": "Please enable JavaScript to continue reading this article."},
            {"entity": "Beta Energy", "notes": "Cookie Settings Privacy Policy Terms of Use"},
        ],
    )
    return "FINAL REPORT:\nCreated CSV deliverable; hygiene should warn on structured crawl contamination."


def _constraint_drift_negative(workspace: Path) -> str:
    report = workspace / "reports" / "date-drift.md"
    report.parent.mkdir(parents=True, exist_ok=True)
    report.write_text(
        textwrap.dedent(
            f"""
            # Date Drift Negative Control

            {PUBLIC_MARKER}

            This report improperly uses a source dated July 2024 even though the request ended at May 2024.
            """
        ).strip()
        + "\n",
        encoding="utf-8",
    )
    return "FINAL REPORT:\nCreated Markdown report."


def _missing_artifact_negative(workspace: Path) -> str:
    report = workspace / "reports" / "notes-only.md"
    report.parent.mkdir(parents=True, exist_ok=True)
    report.write_text(
        f"# Notes Only Negative Control\n\n{PUBLIC_MARKER}\n\nThis intentionally omits requested PDF and XLSX files.\n",
        encoding="utf-8",
    )
    return "FINAL REPORT:\nCreated notes only, but requested artifacts are missing."


def _notes_only_deliverable_intent_negative(workspace: Path) -> str:
    report = workspace / "research" / "batch-01.md"
    report.parent.mkdir(parents=True, exist_ok=True)
    report.write_text(
        f"# Working Notes Only\n\n{PUBLIC_MARKER}\n\nThis intentionally leaves only internal notes.\n",
        encoding="utf-8",
    )
    return "FINAL REPORT:\nCreated internal notes but no final report/table deliverable."


def _planning_window_widen_negative(workspace: Path) -> str:
    specs = workspace / "specs"
    specs.mkdir(parents=True, exist_ok=True)
    (specs / "SPEC.md").write_text(
        "Plan: use public sources from January 2024 through July 2026 for the final screen.\n",
        encoding="utf-8",
    )
    report = workspace / "reports" / "screen.md"
    report.parent.mkdir(parents=True, exist_ok=True)
    report.write_text(
        f"# Planning Drift Negative Control\n\n{PUBLIC_MARKER}\n\nThe visible output is present, but the plan widened the source window.\n",
        encoding="utf-8",
    )
    return "FINAL REPORT:\nCreated Markdown screen."


def _structured_final_marker_negative(workspace: Path) -> str:
    return "Status note: instructions say to end with FINAL REPORT:, but no final answer was captured."


def _xlsx_source_window_negative(workspace: Path) -> str:
    _write_xlsx(
        workspace / "output" / "source-ledger.xlsx",
        [
            ["entity", "source_note"],
            ["Synthetic Target", "Primary source published July 2026 and used for scoring."],
        ],
    )
    return "FINAL REPORT:\nCreated XLSX source ledger."


def _binary_extraction_unavailable_warn(workspace: Path) -> str:
    path = workspace / "output" / "oversized-source-ledger.xlsx"
    path.parent.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_STORED) as archive:
        archive.writestr("[Content_Types].xml", "<Types></Types>")
        archive.writestr("xl/workbook.xml", "<workbook></workbook>")
        archive.writestr("xl/worksheets/sheet1.xml", "<worksheet><sheetData /></worksheet>")
        archive.writestr("xl/media/padding.bin", b"0" * 2_200_000)
    return "FINAL REPORT:\nCreated oversized XLSX source ledger."


def _coverage_range_pass(workspace: Path) -> str:
    rows = [{"entity": f"Synthetic Target {index:02d}", "score": "4.0"} for index in range(12)]
    _write_csv(workspace / "output" / "coverage-screen.csv", rows)
    return "FINAL REPORT:\nCreated coverage screen with 12 synthetic targets."


def _coverage_range_low_negative(workspace: Path) -> str:
    rows = [{"entity": f"Synthetic Target {index:02d}", "score": "4.0"} for index in range(7)]
    _write_csv(workspace / "output" / "coverage-screen.csv", rows)
    return "FINAL REPORT:\nCreated coverage screen with too few synthetic targets."


CASES = [
    BenchmarkCase(
        case_id="generic_market_research_package",
        description="Public-safe market research with seed coverage and Markdown/CSV/XLSX deliverables.",
        instruction=(
            "Use only public official source placeholders dated through June 2026. "
            "Include seed entities Alpha Robotics and Beta Energy. "
            "Deliver Markdown report, CSV screen, and XLSX workbook. Do not create PDF."
        ),
        expected_status="pass",
        expected_constraint_status="pass",
        expected_completion_status="pass",
        expected_coverage_status="not_applicable",
        expected_hygiene_status="pass",
        setup=_market_research_package,
    ),
    BenchmarkCase(
        case_id="technical_literature_brief",
        description="Public-safe technical literature brief with Markdown and JSON evidence.",
        instruction=(
            "Use public technical literature through June 2026. "
            "Include seed topics container sandbox controls and local MCP transport auth. "
            "Deliver Markdown brief and JSON evidence ledger."
        ),
        expected_status="pass",
        expected_constraint_status="pass",
        expected_completion_status="pass",
        expected_coverage_status="not_applicable",
        expected_hygiene_status="pass",
        setup=_technical_literature_brief,
    ),
    BenchmarkCase(
        case_id="file_transform_package",
        description="Uploaded CSV input should not become a required CSV output; professional files validate.",
        instruction=(
            "Using the attached input CSV as the source, create PDF, DOCX, PPTX, and HTML deliverables. "
            "Do not create a CSV output file."
        ),
        expected_status="pass",
        expected_constraint_status="pass",
        expected_completion_status="pass",
        expected_coverage_status="not_applicable",
        expected_hygiene_status="pass",
        setup=_file_transform_package,
    ),
    BenchmarkCase(
        case_id="text_only_no_file_answer",
        description="Text-only answer should pass without forced artifacts.",
        instruction="Answer directly in the final report only. No files or artifacts.",
        expected_status="pass",
        expected_constraint_status="pass",
        expected_completion_status="pass",
        expected_coverage_status="not_applicable",
        expected_hygiene_status="pass",
        setup=_text_only_answer,
    ),
    BenchmarkCase(
        case_id="contaminated_structured_field_warn",
        description="Structured crawl/navigation contamination should warn but remain advisory.",
        instruction="Deliver a CSV research summary using public source placeholders through June 2026.",
        expected_status="warn",
        expected_constraint_status="pass",
        expected_completion_status="pass",
        expected_coverage_status="not_applicable",
        expected_hygiene_status="warn",
        setup=_contaminated_structured_field,
    ),
    BenchmarkCase(
        case_id="constraint_drift_negative_control",
        description="Out-of-window source date should fail constraint compliance.",
        instruction="Use sources only through May 2024. Deliver Markdown report.",
        expected_status="fail",
        expected_constraint_status="fail",
        expected_completion_status="pass",
        expected_coverage_status="not_applicable",
        expected_hygiene_status="pass",
        setup=_constraint_drift_negative,
    ),
    BenchmarkCase(
        case_id="missing_requested_artifact_negative_control",
        description="Missing requested professional artifacts should fail completion compliance.",
        instruction="Deliver PDF and XLSX artifacts.",
        expected_status="fail",
        expected_constraint_status="pass",
        expected_completion_status="fail",
        expected_coverage_status="not_applicable",
        expected_hygiene_status="pass",
        setup=_missing_artifact_negative,
    ),
    BenchmarkCase(
        case_id="notes_only_deliverable_intent_negative_control",
        description="Internal notes only should fail when the request asks for report/table deliverables without explicit extensions.",
        instruction=(
            "### OUTPUT FORMAT\n"
            "Section 1 - Master Summary Table. All targets, sortable by weighted fit score.\n"
            "Section 2 - Tier 1 Deep Profiles. Deliver the final report and recommendation."
        ),
        expected_status="fail",
        expected_constraint_status="pass",
        expected_completion_status="fail",
        expected_coverage_status="not_applicable",
        expected_hygiene_status="pass",
        setup=_notes_only_deliverable_intent_negative,
    ),
    BenchmarkCase(
        case_id="planning_source_window_widen_negative_control",
        description="Planning/spec files that widen a strict source/date window should fail constraint compliance.",
        instruction="Use public sources from January 2024 through May 2026 only. Deliver a Markdown screen.",
        expected_status="fail",
        expected_constraint_status="fail",
        expected_completion_status="pass",
        expected_coverage_status="not_applicable",
        expected_hygiene_status="pass",
        setup=_planning_window_widen_negative,
    ),
    BenchmarkCase(
        case_id="structured_final_marker_negative_control",
        description="Progress text that merely mentions the required final marker should not count as a final report.",
        instruction="Answer directly with a final report.",
        expected_status="fail",
        expected_constraint_status="pass",
        expected_completion_status="fail",
        expected_coverage_status="not_applicable",
        expected_hygiene_status="pass",
        setup=_structured_final_marker_negative,
    ),
    BenchmarkCase(
        case_id="xlsx_source_window_negative_control",
        description="Out-of-window source evidence hidden in a generated workbook should fail constraint compliance.",
        instruction="Use public sources from January 2024 through May 2026 only. Deliver XLSX source ledger.",
        expected_status="fail",
        expected_constraint_status="fail",
        expected_completion_status="pass",
        expected_coverage_status="not_applicable",
        expected_hygiene_status="pass",
        setup=_xlsx_source_window_negative,
    ),
    BenchmarkCase(
        case_id="binary_extraction_unavailable_warn_control",
        description="Binary artifact text extraction gaps should warn instead of silently passing.",
        instruction="Use public sources from January 2024 through May 2026 only. Deliver XLSX source ledger.",
        expected_status="warn",
        expected_constraint_status="warn",
        expected_completion_status="pass",
        expected_coverage_status="not_applicable",
        expected_hygiene_status="pass",
        setup=_binary_extraction_unavailable_warn,
    ),
    BenchmarkCase(
        case_id="coverage_range_pass_control",
        description="Requested coverage count range should pass when deliverable table rows fall inside the range.",
        instruction="Aim for 10-15 targets total before scoring. Deliver CSV screen.",
        expected_status="pass",
        expected_constraint_status="pass",
        expected_completion_status="pass",
        expected_coverage_status="pass",
        expected_hygiene_status="pass",
        setup=_coverage_range_pass,
    ),
    BenchmarkCase(
        case_id="coverage_range_low_negative_control",
        description="Requested coverage count range should fail when deliverable table rows are below the minimum.",
        instruction="Aim for 10-15 targets total before scoring. Deliver CSV screen.",
        expected_status="fail",
        expected_constraint_status="pass",
        expected_completion_status="pass",
        expected_coverage_status="fail",
        expected_hygiene_status="pass",
        setup=_coverage_range_low_negative,
    ),
]


def _run_case(case: BenchmarkCase, state_root: Path) -> dict[str, object]:
    workspace = state_root / case.case_id
    if workspace.exists():
        shutil.rmtree(workspace)
    workspace.mkdir(parents=True)
    output_text = case.setup(workspace)
    run_id = f"run_{case.case_id}"
    worker = {
        "worker_id": f"wrk_{case.case_id}",
        "profile": "codex-cli",
        "execution_mode": "host",
        "runtime": "codex-cli",
        "backend": "openclaw",
    }
    ledger = build_constraint_ledger(instruction=case.instruction, worker=worker, run_id=run_id)
    write_constraint_ledger(workspace, ledger, run_id)
    evidence = build_run_evidence(
        worker=worker,
        run_id=run_id,
        runtime_name="codex-cli",
        model="public-safe-fixture",
        command=["codex", "exec", "-c", 'model_reasoning_effort="xhigh"', "-"],
        env={},
        workspace_dir=workspace,
        stdout_text=output_text,
        stderr_text="",
        output_text=output_text,
        error_text="",
        exit_code=0,
        timeout_seconds=300,
        stop_reason="process_exit",
        constraint_ledger=ledger,
        started_at=time.time(),
        ended_at=time.time(),
        transcript_paths={"stdout_tail": "glasshive-run/runs/{run_id}/stdout-tail.txt"},
    )
    write_run_evidence(workspace, evidence, run_id)
    actual = {
        "evidence_status": str(evidence["evidence_result"]["status"]),
        "constraint_status": str(evidence["constraint_compliance"]["status"]),
        "completion_status": str(evidence["completion_compliance"]["status"]),
        "coverage_status": str(evidence["coverage_compliance"]["status"]),
        "coverage_observed_max_count": int(evidence["coverage_compliance"].get("observed_max_count") or 0),
        "hygiene_status": str(evidence["content_hygiene"]["status"]),
        "artifact_count": int(evidence["artifacts"]["count"]),
        "required_artifact_types": list(evidence["completion_compliance"].get("required_artifact_types", [])),
        "missing_required_artifact_types": list(
            evidence["completion_compliance"].get("missing_required_artifact_types", [])
        ),
    }
    expected = {
        "evidence_status": case.expected_status,
        "constraint_status": case.expected_constraint_status,
        "completion_status": case.expected_completion_status,
        "coverage_status": case.expected_coverage_status,
        "hygiene_status": case.expected_hygiene_status,
    }
    return {
        "case_id": case.case_id,
        "description": case.description,
        "expected": expected,
        "actual": actual,
        "passed": all(actual[key] == expected[key] for key in expected),
    }


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--state-root",
        type=Path,
        default=Path(tempfile.gettempdir()) / "glasshive-public-safe-benchmarks",
    )
    parser.add_argument("--output", type=Path, default=None, help="Optional JSON summary output path.")
    args = parser.parse_args()

    args.state_root.mkdir(parents=True, exist_ok=True)
    results = [_run_case(case, args.state_root) for case in CASES]
    summary = {
        "schema": "glasshive.public-safe-benchmark-summary.v1",
        "case_count": len(results),
        "passed_count": sum(1 for item in results if item["passed"]),
        "failed_count": sum(1 for item in results if not item["passed"]),
        "results": results,
    }
    body = json.dumps(summary, indent=2, sort_keys=True)
    if args.output:
        args.output.parent.mkdir(parents=True, exist_ok=True)
        args.output.write_text(body + "\n", encoding="utf-8")
    print(body)
    return 0 if summary["failed_count"] == 0 else 1


if __name__ == "__main__":
    raise SystemExit(main())
